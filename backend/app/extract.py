from __future__ import annotations

import io
import re
from pathlib import Path

from pypdf import PdfReader
from pptx import Presentation


def extract_text(filename: str, data: bytes) -> str:
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        reader = PdfReader(io.BytesIO(data))
        pages = []
        for page in reader.pages:
            pages.append(page.extract_text() or "")
        return "\n".join(pages)
    if ext in {".ppt", ".pptx"}:
        prs = Presentation(io.BytesIO(data))
        bits = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    bits.append(shape.text)
        return "\n".join(bits)
    if ext in {".txt", ".md", ".csv"}:
        return data.decode("utf-8", errors="ignore")
    return data.decode("utf-8", errors="ignore")


def split_concepts(text: str, fallback_title: str) -> list[dict]:
    cleaned = re.sub(r"\r", "", text or "").strip()
    if len(cleaned) < 80:
        return _aptitude_seed(fallback_title)

    chunks: list[str] = []
    heading_split = re.split(r"\n(?=[A-Z][A-Za-z0-9 .:?-]{8,80}\n)", cleaned)
    if len(heading_split) < 3:
        heading_split = re.split(r"\n{2,}", cleaned)
    for part in heading_split:
        part = part.strip()
        if len(part) < 40:
            continue
        chunks.append(part)
        if len(chunks) >= 10:
            break
    if not chunks:
        return _aptitude_seed(fallback_title)

    concepts = []
    for i, chunk in enumerate(chunks):
        title = chunk.split("\n", 1)[0][:80].strip() or f"Concept {i + 1}"
        body = chunk[:1800]
        difficulty = "easy" if i == 0 else "medium" if i < 4 else "hard"
        concepts.append(
            {
                "title": title,
                "body": body,
                "difficulty": difficulty,
                "depends_on": chunks[i - 1].split("\n", 1)[0][:80] if i else "",
            }
        )
    return concepts


def _aptitude_seed(title: str) -> list[dict]:
    samples = [
        (
            "Percentages",
            "A percentage is a way of expressing a number as a fraction of 100. If 25 out of 200 students pass, that is 12.5%. Increase and decrease questions often hide a two-step change.",
            "easy",
        ),
        (
            "Ratios and proportions",
            "A ratio compares quantities of the same kind. If A:B = 2:3 and B:C = 4:5, make B the same in both before combining. Proportions keep two ratios equal.",
            "easy",
        ),
        (
            "Time, speed and distance",
            "Distance = speed × time. Average speed is total distance over total time, not the average of the two speeds. Upstream/downstream and relative speed are the usual twists.",
            "medium",
        ),
        (
            "Profit and loss",
            "Profit % is always on cost price unless the question says otherwise. Marked price, discount, and successive discounts are common exam traps.",
            "medium",
        ),
        (
            "Permutations and combinations",
            "Permutations care about order; combinations do not. nPr = n! / (n-r)! and nCr = n! / (r!(n-r)!). Word problems hide whether order matters.",
            "hard",
        ),
        (
            "Data interpretation",
            "Read the chart title and units first. Compare, don't just read off a number. Approximation saves time when options are far apart.",
            "hard",
        ),
    ]
    out = []
    for i, (t, body, d) in enumerate(samples):
        out.append(
            {
                "title": t if i else f"{title}: {t}",
                "body": body,
                "difficulty": d,
                "depends_on": samples[i - 1][0] if i else "",
            }
        )
    return out
