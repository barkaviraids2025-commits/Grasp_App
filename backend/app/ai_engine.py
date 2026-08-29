import json
import os
import re
from pathlib import Path
from dotenv import load_dotenv

load_dotenv()

from pypdf import PdfReader

try:
    from pptx import Presentation
except Exception:  # pragma: no cover
    Presentation = None


APTITUDE_HINTS = [
    ("percentages", "Percentages", "medium"),
    ("percent", "Percentages", "medium"),
    ("ratio", "Ratios and Proportion", "medium"),
    ("proportion", "Ratios and Proportion", "medium"),
    ("average", "Averages", "easy"),
    ("profit", "Profit and Loss", "medium"),
    ("loss", "Profit and Loss", "medium"),
    ("discount", "Profit, Loss and Discount", "medium"),
    ("time and work", "Time and Work", "hard"),
    ("work and wages", "Time and Work", "hard"),
    ("pipes", "Pipes and Cisterns", "medium"),
    ("cistern", "Pipes and Cisterns", "medium"),
    ("time, speed", "Time, Speed and Distance", "hard"),
    ("speed", "Time, Speed and Distance", "hard"),
    ("train", "Problems on Trains", "medium"),
    ("boat", "Boats and Streams", "hard"),
    ("stream", "Boats and Streams", "hard"),
    ("permutation", "Permutations and Combinations", "hard"),
    ("combination", "Permutations and Combinations", "hard"),
    ("probability", "Probability", "hard"),
    ("number system", "Number System", "easy"),
    ("divisibility", "Number System & Divisibility", "easy"),
    ("hcf", "HCF and LCM", "easy"),
    ("lcm", "HCF and LCM", "easy"),
    ("algebra", "Algebraic Equations", "medium"),
    ("quadratic", "Quadratic Equations", "hard"),
    ("geometry", "Geometry Fundamentals", "medium"),
    ("triangle", "Triangles & Polygons", "medium"),
    ("circle", "Circles & Tangents", "medium"),
    ("mensuration", "Mensuration (2D & 3D Area/Volume)", "medium"),
    ("volume", "Mensuration (Area & Volume)", "medium"),
    ("data interpretation", "Data Interpretation & Charts", "medium"),
    ("pie chart", "Data Interpretation (Pie Charts)", "medium"),
    ("bar chart", "Data Interpretation (Bar Graphs)", "medium"),
    ("simple interest", "Simple Interest (SI)", "medium"),
    ("compound interest", "Compound Interest (CI)", "medium"),
    ("interest", "Simple and Compound Interest", "medium"),
    ("ages", "Problems on Ages", "easy"),
    ("mixture", "Mixtures and Alligations", "hard"),
    ("alligation", "Mixtures and Alligations", "hard"),
    ("series", "Number & Letter Series", "easy"),
    ("coding", "Coding-Decoding Reasoning", "easy"),
    ("blood relation", "Blood Relations", "easy"),
    ("syllogism", "Syllogisms & Deductions", "medium"),
    ("seating", "Seating Arrangement", "hard"),
    ("clock", "Clocks & Calendars", "medium"),
    ("calendar", "Clocks & Calendars", "medium"),
]

EXPLANATION_MODES = ["simple", "example", "diagram", "animation", "voice", "practice"]


def next_mode(current: str) -> str:
    if current not in EXPLANATION_MODES:
        return "example"
    idx = (EXPLANATION_MODES.index(current) + 1) % len(EXPLANATION_MODES)
    return EXPLANATION_MODES[idx]


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        try:
            reader = PdfReader(str(path))
            return "\n".join((page.extract_text() or "") for page in reader.pages)
        except Exception:
            return ""
    if suffix in {".pptx", ".ppt"} and Presentation:
        try:
            pres = Presentation(str(path))
            chunks = []
            for slide in pres.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        chunks.append(shape.text)
            return "\n".join(chunks)
        except Exception:
            return ""
    if suffix in {".txt", ".md", ".csv"}:
        return path.read_text(encoding="utf-8", errors="ignore")
    try:
        return path.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return ""


def _split_blocks(text: str) -> list[str]:
    cleaned = re.sub(r"\r", "", text)
    parts = re.split(r"\n\s*\n|Chapter\s+\d+|Unit\s+\d+|Topic\s+\d+|Module\s+\d+", cleaned, flags=re.I)
    blocks = [re.sub(r"\s+", " ", p).strip() for p in parts if len(p.strip()) > 35]
    if len(blocks) < 3:
        sentences = re.split(r"(?<=[.!?])\s+", cleaned)
        chunk, acc = [], []
        for s in sentences:
            acc.append(s)
            if len(" ".join(acc)) > 260:
                chunk.append(" ".join(acc).strip())
                acc = []
        if acc:
            chunk.append(" ".join(acc).strip())
        blocks = [c for c in chunk if len(c) > 35]
    return blocks[:18]


def _title_for(block: str, index: int) -> tuple[str, str]:
    lower = block.lower()
    for needle, title, difficulty in APTITUDE_HINTS:
        if needle in lower:
            return title, difficulty
    words = [w for w in block.split() if len(w) > 2][:6]
    title = " ".join(words).strip(" :-.,;")
    if len(title) < 5:
        title = f"Aptitude Concept {index + 1}"
    return title.title()[:70], "medium"


def analyze_source(path: Path) -> list[dict]:
    text = extract_text(path)
    if not text.strip():
        return [
            {
                "title": "Quantitative Aptitude Foundations",
                "summary": "Core quantitative analysis fundamentals, formulas, and problem-solving shortcuts.",
                "source_excerpt": "Comprehensive quantitative aptitude principles structured for fast understanding and systematic practice.",
                "difficulty": "medium",
            },
            {
                "title": "Ratios, Percentages & Averages",
                "summary": "Key baseline arithmetic concepts and shortcut calculation techniques.",
                "source_excerpt": "Percentage calculations, ratio simplification, and weighted averages applied to exam problems.",
                "difficulty": "medium",
            },
            {
                "title": "Time, Speed, Distance & Work",
                "summary": "Motion formulas, relative speed, and inverse-proportion work rates.",
                "source_excerpt": "Standard problems on trains, work-time equations, and rate combinations.",
                "difficulty": "hard",
            },
            {
                "title": "Data Interpretation & Problem Solving",
                "summary": "Reading charts, tables, and logical extraction under timed conditions.",
                "source_excerpt": "Interpreting numerical data quickly and accurately without long calculation.",
                "difficulty": "medium",
            },
        ]
    concepts = []
    seen = set()
    for i, block in enumerate(_split_blocks(text)):
        title, difficulty = _title_for(block, i)
        key = title.lower()
        if key in seen:
            title = f"{title} (Part {i + 1})"
        seen.add(title.lower())
        excerpt = block[:900]
        concepts.append(
            {
                "title": title,
                "summary": _local_simple(title, excerpt),
                "source_excerpt": excerpt,
                "difficulty": difficulty,
            }
        )
    return concepts[:12] or [
        {
            "title": "Practice Set from Source",
            "summary": "Step-by-step concepts extracted from your uploaded material.",
            "source_excerpt": text[:900],
            "difficulty": "medium",
        }
    ]


def _local_simple(title: str, excerpt: str) -> str:
    return (
        f"{title} is a core quantitative aptitude concept. "
        f"Master the fundamental relationship first, visualize with an everyday analogy, "
        f"then lock in the formula with practice. Source context: {excerpt[:200]}..."
    )


def explanations_for(concept: dict, language: str, mode: str) -> dict:
    title = concept.get("title", "Concept")
    excerpt = concept.get("source_excerpt", "")
    is_tamil = language in ("ta", "tamil")

    # Local structured fallback templates
    simple_en = (
        f"**{title} in plain English:**\n"
        f"Quantitative analysis isn't about memorizing massive formulas—it's about naming what you HAVE (given values), "
        f"what you NEED (target), and the single relationship linking them.\n\n"
        f"💡 **Core Rule:** Always standardize units first, write down the formula, and estimate the magnitude "
        f"before doing heavy calculations."
    )
    simple_ta = (
        f"**{title} - எளிய விளக்கம்:**\n"
        f"{title} என்பது கடினமான சூத்திரங்களை மனப்பாடம் செய்வதல்ல; கொடுக்கப்பட்ட அளவுகள் (Given), "
        f"நாம் கண்டுபிடிக்க வேண்டியது (Target) ஆகியவற்றை சரியாகப் புரிந்து கொள்வதே ஆகும்.\n\n"
        f"💡 **முக்கிய விதி:** முதலில் அலகுகளை (units) ஒரே வகையாக மாற்றவும், பின்னர் சூத்திரத்தைப் பயன்படுத்தி எளிய முறையில் கணக்கிடவும்."
    )

    example_en = (
        f"**Real-world scenario for {title}:**\n"
        f"Imagine you are running a small business or buying groceries. "
        f"If a store offers a 20% discount on an item marked at $150, you pay 80% of $150 = $120. "
        f"The same ratio and scaling logic applies to any exam question on {title}."
    )
    example_ta = (
        f"**தினசரி வாழ்க்கை உதாரணம் ({title}):**\n"
        f"ஒரு கடையில் ₹100 மதிப்புள்ள பொருளுக்கு 20% தள்ளுபடி கிடைத்தால், நாம் செலுத்தும் தொகை ₹80. "
        f"இதே எளிய சதவீத மற்றும் விகித முறையே {title} தொடர்பான அனைத்து வினாக்களுக்கும் அடிப்படையாகும்."
    )

    diagram_en = (
        f"**Visual Mental Diagram for {title}:**\n"
        f"┌─────────────────────────┐      ┌─────────────────────────┐\n"
        f"│   GIVEN / INPUTS        │ ───► │   CONNECTING FORMULA    │\n"
        f"│ (Numbers, Units, Rates) │      │      ({title})          │\n"
        f"└─────────────────────────┘      └───────────┬─────────────┘\n"
        f"                                             │\n"
        f"                                             ▼\n"
        f"                                 ┌─────────────────────────┐\n"
        f"                                 │  TARGET ANSWER + CHECK  │\n"
        f"                                 │  (Does magnitude fit?)  │\n"
        f"                                 └─────────────────────────┘"
    )
    diagram_ta = (
        f"**மன வரைபடம் (Mental Diagram - {title}):**\n"
        f"[ கொடுக்கப்பட்ட மதிப்புகள் ] ───► [ {title} சூத்திரம் ] ───► [ சரியான தீர்வு ]\n"
        f"எப்போதும் இறுதி விடையின் அளவு இயல்புக்கு ஒத்துப்போகிறதா என்று சரிபார்க்கவும்."
    )

    animation_en = (
        f"**Process Animation / Step-by-Step Flow:**\n"
        f"1. 📥 **Step 1:** Extract the raw numbers from the problem statement.\n"
        f"2. 🔄 **Step 2:** Convert all quantities to uniform units (e.g. km/h to m/s, or hours to minutes).\n"
        f"3. ⚡ **Step 3:** Apply the direct ratio/formula for {title}.\n"
        f"4. 🎯 **Step 4:** Cancel out common factors before multiplying to save time."
    )
    animation_ta = (
        f"**படிமுறை ஓட்டம் (Step-by-step Process):**\n"
        f"1. 📥 **படி 1:** வினாவில் உள்ள எண்களைத் தனித்தனியாகப் பிரிக்கவும்.\n"
        f"2. 🔄 **படி 2:** அலகுகளை ஒரே மாதிரியாக மாற்றவும்.\n"
        f"3. ⚡ **படி 3:** {title} நேரடி விதியைப் பயன்படுத்தவும்.\n"
        f"4. 🎯 **படி 4:** தேவையற்ற எண்களை சுருக்கி விடையைக் காணவும்."
    )

    practice_en = (
        f"**Quick Insight & Shortcut:**\n"
        f"Whenever you see a {title} problem with multiple options, test the extremes or use elimination. "
        f"If the value must increase, instantly cross off any option lower than the starting number!"
    )
    practice_ta = (
        f"**குறுக்குவழி உத்தி (Shortcut Tip):**\n"
        f"{title} வினாக்களில் விடைகளைத் தேர்வு செய்யும் போது, தவறான விடைகளை (Elimination Method) "
        f"முதலில் நீக்குவது மிக வேகமாக விடை காண உதவும்."
    )

    voice_en = (
        f"Let us walk through {title} slowly and clearly. "
        f"First, identify what is given. Second, identify what is asked. "
        f"Third, apply the core relationship. Finally, do a 2-second sanity check. You've got this!"
    )
    voice_ta = (
        f"{title} தலைப்பை மிக அமைதியாகப் புரிந்து கொள்வோம். "
        f"முதலில் தரப்பட்ட எண்களைக் கவனியுங்கள். அடுத்து கேட்கப்பட்ட கேள்வியைக் கவனியுங்கள். "
        f"இவ்விரண்டையும் இணைக்கும் விதியை அமைத்தால் விடை மிக எளிதாக வரும்!"
    )

    payload = {
        "simple": simple_ta if is_tamil else simple_en,
        "example": example_ta if is_tamil else example_en,
        "diagram": diagram_ta if is_tamil else diagram_en,
        "animation": animation_ta if is_tamil else animation_en,
        "practice": practice_ta if is_tamil else practice_en,
        "voice": voice_ta if is_tamil else voice_en,
        "grounded": excerpt[:400],
        "mode": mode,
        "title": title,
    }

    # Try Gemini or OpenAI if configured
    gemini_key = os.getenv("GEMINI_API_KEY") or os.getenv("GOOGLE_API_KEY")
    openai_key = os.getenv("OPENAI_API_KEY")

    if gemini_key:
        generated = _gemini_explain(gemini_key, concept, language, mode)
        if generated:
            payload.update(generated)
    elif openai_key:
        generated = _openai_explain(openai_key, concept, language, mode)
        if generated:
            payload.update(generated)

    return payload


def _gemini_explain(api_key: str, concept: dict, language: str, mode: str) -> dict | None:
    try:
        import httpx

        lang_name = "Tamil (native, natural language re-explanation with English mathematical terms)" if language in ("ta", "tamil") else "English"
        prompt = (
            f"You are Concepta, an expert adaptive Aptitude and Quantitative Analysis personal coach.\n"
            f"Teach the concept '{concept.get('title')}' in an intuitive, engaging way.\n"
            f"Language: {lang_name}.\n"
            f"Current Focus Mode: {mode}.\n"
            f"Source material excerpt: {concept.get('source_excerpt', '')[:1000]}\n\n"
            f"Output ONLY valid JSON with exactly these keys:\n"
            f'{{\n'
            f'  "simple": "Clear, intuitive breakdown without jargon in {lang_name}",\n'
            f'  "example": "Relatable everyday real-world example with numbers in {lang_name}",\n'
            f'  "diagram": "Text-based ASCII or mental model diagram explaining the formula structure",\n'
            f'  "animation": "Step-by-step visual animation / logical flow (Step 1, Step 2, Step 3)",\n'
            f'  "practice": "Pro-tip / shortcut calculation trick for competitive exams",\n'
            f'  "voice": "A friendly, encouraging spoken-style walkthrough script in {lang_name}"\n'
            f'}}'
        )

        url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent?key={api_key}"
        resp = httpx.post(
            url,
            headers={"Content-Type": "application/json"},
            json={
                "contents": [{"parts": [{"text": prompt}]}],
                "generationConfig": {"temperature": 0.4, "maxOutputTokens": 1200},
            },
            timeout=20,
        )
        if resp.status_code >= 400:
            # Try gemini-2.5-flash or gemini-1.5-pro fallback
            url2 = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent?key={api_key}"
            resp = httpx.post(
                url2,
                headers={"Content-Type": "application/json"},
                json={
                    "contents": [{"parts": [{"text": prompt}]}],
                    "generationConfig": {"temperature": 0.4, "maxOutputTokens": 1200},
                },
                timeout=20,
            )
            if resp.status_code >= 400:
                return None

        data = resp.json()
        text_content = data["candidates"][0]["content"]["parts"][0]["text"]
        start, end = text_content.find("{"), text_content.rfind("}")
        if start >= 0 and end > start:
            return json.loads(text_content[start : end + 1])
    except Exception:
        return None
    return None


def _openai_explain(api_key: str, concept: dict, language: str, mode: str) -> dict | None:
    try:
        import httpx

        lang = "Tamil" if language in ("ta", "tamil") else "English"
        prompt = (
            f"You are Concepta, an aptitude learning coach. Re-explain, do not merely translate.\n"
            f"Language: {lang}. Mode: {mode}. Concept: {concept['title']}.\n"
            f"Source excerpt: {concept.get('source_excerpt', '')[:1000]}.\n"
            "Return JSON with keys: simple, example, diagram, animation, practice, voice."
        )
        resp = httpx.post(
            "https://api.openai.com/v1/chat/completions",
            headers={"Authorization": f"Bearer {api_key}"},
            json={
                "model": os.getenv("OPENAI_MODEL", "gpt-4o-mini"),
                "messages": [{"role": "user", "content": prompt}],
                "temperature": 0.4,
            },
            timeout=20,
        )
        if resp.status_code >= 400:
            return None
        content = resp.json()["choices"][0]["message"]["content"]
        start, end = content.find("{"), content.rfind("}")
        if start >= 0 and end > start:
            return json.loads(content[start : end + 1])
    except Exception:
        return None
    return None


def generate_questions(concept: dict, attempt: int = 0) -> list[dict]:
    title = concept.get("title", "Aptitude Concept")
    harder = attempt >= 1
    
    return [
        {
            "id": "mcq1",
            "kind": "mcq",
            "prompt": f"What is the most effective first move when starting a problem on '{title}'?"
            if not harder
            else f"When tackling a complex or unfamiliar '{title}' problem, what ensures accuracy?",
            "options": [
                "List known inputs, target variable, and ensure uniform measurement units",
                "Immediately perform multiplication with all numbers given",
                "Rely on memorized final answer patterns without checking units",
                "Guess between the two highest numerical choices",
            ],
            "answer": 0,
        },
        {
            "id": "scenario",
            "kind": "mcq",
            "prompt": f"In a realistic test question involving '{title}', values are given in different dimensions (e.g. hours vs minutes or km vs m). What is crucial?",
            "options": [
                "Convert all quantities to consistent baseline units before computing",
                "Multiply directly and convert units only at the end if you have time",
                "Ignore unit differences because ratios cancel them out automatically",
                "Add the numbers in the order they are presented in the question",
            ],
            "answer": 0,
        },
        {
            "id": "own_words",
            "kind": "own_words",
            "prompt": f"Explain '{title}' in your own words. Include one small example or shortcut rule to prove real understanding.",
        },
    ]


def score_own_words(text: str, title: str) -> float:
    if not text or not text.strip():
        return 0.0
    words = re.findall(r"[a-zA-Z\u0B80-\u0BFF]+", text.lower())
    if len(words) < 5:
        return 0.2
    tokens = set(re.findall(r"[a-z0-9\u0B80-\u0BFF]+", title.lower()))
    overlap = len(tokens & set(words))
    length_bonus = min(len(words) / 35, 0.45)
    return min(1.0, 0.4 + 0.15 * overlap + length_bonus)
